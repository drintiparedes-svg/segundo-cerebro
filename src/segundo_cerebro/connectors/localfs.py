"""Conector de carpetas locales — SOLO LECTURA.

Registra carpetas de la máquina (por ejemplo las del escritorio) como
fuentes de información. Garantía de diseño: este módulo abre los archivos
exclusivamente en modo lectura; nunca crea, modifica ni borra nada dentro
de una fuente. Todo lo que el sistema escribe vive en .brain/.

Tipos soportados: .md, .txt, .csv, .html, .json siempre; .pdf, .docx,
.xlsx y .pptx si están instalados los extras (pip install -e ".[files]").
"""

from __future__ import annotations

import json
from datetime import datetime, timezone
from pathlib import Path

from ..models import Document, new_id, now_iso

TEXT_EXTS = {".md", ".txt", ".csv", ".json"}
OPTIONAL_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".html", ".htm"}
SKIP_DIRS = {".git", ".brain", "node_modules", "__pycache__", ".venv",
             "$RECYCLE.BIN", "System Volume Information"}
SKIP_FILES = {"desktop.ini", "Thumbs.db", ".DS_Store", ".localized"}
MAX_FILE_BYTES = 15_000_000
MAX_BODY_CHARS = 200_000


# ── registro de fuentes (.brain/sources.json) ─────────────────────────────

def _registry_path(brain_dir: str | Path) -> Path:
    p = Path(brain_dir)
    p.mkdir(parents=True, exist_ok=True)
    return p / "sources.json"


def load_registry(brain_dir: str | Path) -> dict:
    path = _registry_path(brain_dir)
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {"sources": [], "state": {}}


def save_registry(brain_dir: str | Path, registry: dict) -> None:
    _registry_path(brain_dir).write_text(
        json.dumps(registry, ensure_ascii=False, indent=2), encoding="utf-8")


def add_source(brain_dir: str | Path, folder: str | Path,
               alias: str | None = None) -> dict:
    folder = Path(folder).expanduser().resolve()
    if not folder.is_dir():
        raise NotADirectoryError(f"No es una carpeta: {folder}")
    registry = load_registry(brain_dir)
    entry = {"path": str(folder), "alias": alias or folder.name,
             "added_at": now_iso()}
    if any(s["path"] == entry["path"] for s in registry["sources"]):
        return entry  # ya registrada
    registry["sources"].append(entry)
    save_registry(brain_dir, registry)
    return entry


def remove_source(brain_dir: str | Path, folder: str | Path) -> bool:
    folder = str(Path(folder).expanduser().resolve())
    registry = load_registry(brain_dir)
    before = len(registry["sources"])
    registry["sources"] = [s for s in registry["sources"] if s["path"] != folder]
    registry["state"].pop(folder, None)
    save_registry(brain_dir, registry)
    return len(registry["sources"]) < before


# ── lectura de archivos (nunca escritura) ─────────────────────────────────

def read_file_text(path: Path) -> str | None:
    """Extrae texto de un archivo soportado. Devuelve None si no se puede
    (tipo no soportado, binario, o falta la librería opcional)."""
    ext = path.suffix.lower()
    if ext in TEXT_EXTS:
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except OSError:
            return None
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return None
        try:
            reader = PdfReader(str(path))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception:
            return None
    if ext == ".docx":
        try:
            import docx
        except ImportError:
            return None
        try:
            d = docx.Document(str(path))
            return "\n".join(p.text for p in d.paragraphs)
        except Exception:
            return None
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError:
            return None
        try:
            wb = load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"## Hoja: {ws.title}")
                for row in ws.iter_rows(max_rows=500, values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        lines.append(" | ".join(cells))
            wb.close()
            return "\n".join(lines)
        except Exception:
            return None
    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return None
        try:
            prs = Presentation(str(path))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## Diapositiva {i}")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)
        except Exception:
            return None
    if ext in {".html", ".htm"}:
        return _html_to_text(path)
    return None


def _html_to_text(path: Path) -> str | None:
    from html.parser import HTMLParser

    class Extractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.chunks: list[str] = []
            self._skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip += 1

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self._skip:
                self._skip -= 1

        def handle_data(self, data):
            if not self._skip and data.strip():
                self.chunks.append(data.strip())

    try:
        parser = Extractor()
        parser.feed(path.read_text(encoding="utf-8", errors="replace"))
        return "\n".join(parser.chunks)
    except Exception:
        return None


def file_to_document(path: Path, source_alias: str, body: str) -> Document:
    mtime = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc)
    return Document(
        id=new_id("doc"),
        path=str(path),
        title=path.stem,
        doc_type="note",
        date=mtime.date().isoformat(),
        body=body[:MAX_BODY_CHARS],
        metadata={
            "source": "local-folder",
            "source_alias": source_alias,
            "extension": path.suffix.lower(),
            "modified_time": mtime.isoformat(),
        },
    )


def iter_files(root: Path):
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        if any(part in SKIP_DIRS or part.startswith(".") for part in path.parts):
            continue
        if path.name in SKIP_FILES or path.name.startswith("~$"):
            continue
        if path.suffix.lower() not in TEXT_EXTS | OPTIONAL_EXTS:
            continue
        try:
            if path.stat().st_size > MAX_FILE_BYTES:
                continue
        except OSError:
            continue
        yield path


def sync_source(store, brain_dir: str | Path, source: dict) -> dict:
    """Sincroniza una fuente. Incremental por mtime + deduplicación por
    hash de contenido: correrlo mil veces no duplica nada."""
    root = Path(source["path"])
    result = {"added": 0, "unchanged": 0, "unsupported": 0, "docs": []}
    if not root.is_dir():
        result["error"] = f"carpeta no disponible: {root}"
        return result

    registry = load_registry(brain_dir)
    state = registry["state"].setdefault(source["path"], {})
    last_mtime = state.get("last_mtime", 0.0)
    max_mtime = last_mtime

    for path in iter_files(root):
        mtime = path.stat().st_mtime
        max_mtime = max(max_mtime, mtime)
        if mtime <= last_mtime:
            result["unchanged"] += 1
            continue
        body = read_file_text(path)
        if body is None or not body.strip():
            result["unsupported"] += 1
            continue
        doc = file_to_document(path, source["alias"], body)
        if store.add_document(doc):
            result["added"] += 1
            result["docs"].append(doc)
        else:
            result["unchanged"] += 1

    state["last_mtime"] = max_mtime
    state["last_sync"] = now_iso()
    save_registry(brain_dir, registry)
    return result
